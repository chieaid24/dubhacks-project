import json
import os
from PyPDF2 import PdfReader
from pptx import Presentation

def parse_file(file_path: str):
    # step 1: extract slide text
    if file_path.endswith(".pptx") or file_path.endswith(".pdf"):
        slides_data = file_to_text(file_path)
    else:
        raise ValueError("Unsupported file type")

    return slides_data

def file_to_text(file_path: str):
    pages = []
    if file_path.endswith(".pptx"):
        # Treat PowerPoint slides as pages, same format as PDFs.
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text_chunks = [
                shape.text.strip()
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            ]
            pages.append({
                "page_number": i + 1,
                "content": "\n".join(text_chunks)
            })
    else:
        reader = PdfReader(file_path)
        for i, page in enumerate(reader.pages):
            text = page.extract_text()
            pages.append({
                "page_number": i + 1,
                "content": text.strip() if text else ""
            })
    return pages