from pptx import Presentation
from PyPDF2 import PdfReader
import fitz
import subprocess
import json
import os


def parse_file(file_path: str):
    # Main entry point — detects file type and routes to parser.
    base_name = os.path.splitext(os.path.basename(file_path))[0]

    if file_path.endswith(".pptx") or file_path.endswith(".pdf"):
        image_paths = file_to_images(file_path)
    else:
        raise ValueError("Unsupported file type")

    # --- Step 2: Extract slide text ---
    if file_path.endswith(".pptx"):
        slides_data = parse_pptx_as_pages(file_path)
    elif file_path.endswith(".pdf"):
        slides_data = parse_pdf(file_path)
    else:
        slides_data = []

    # --- Step 3: Save parsed output (optional for debugging) ---
    parsed_output_dir = os.path.join("outputs", "parsed")
    os.makedirs(parsed_output_dir, exist_ok=True)
    output_json_path = os.path.join(parsed_output_dir, f"{base_name}_parsed.json")
    save_parsed_output(slides_data, output_json_path)

    # --- Step 4: Return structured results ---
    return slides_data, image_paths


def parse_pptx_as_pages(file_path: str):
    # Treat PowerPoint slides as pages, same format as PDFs.
    prs = Presentation(file_path)
    pages = []
    for i, slide in enumerate(prs.slides):
        text_chunks = [
            shape.text.strip()
            for shape in slide.shapes
            if hasattr(shape, "text") and shape.text.strip()
        ]
        pages.append({
            "page_number": i + 1,
            "content": "\n".join(text_chunks)
        })
    return pages


def parse_pdf(file_path: str):
    reader = PdfReader(file_path)
    pages = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        pages.append({
            "page_number": i + 1,
            "content": text.strip() if text else ""
        })
    return pages

def file_to_images(file_path: str):
    output_dir = os.path.dirname(file_path)
    if file_path.endswith(".pptx"):
        subprocess.run([
            "libreoffice", "--headless", "--convert-to", "pdf", file_path,
            "--outdir", output_dir
        ])
        file_path = os.path.join(output_dir, os.path.splitext(os.path.basename(file_path))[0] + ".pdf")
    os.makedirs(output_dir, exist_ok=True)
    doc = fitz.open(file_path)
    image_paths = []

    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # high-resolution
        image_path = os.path.join(output_dir, f"page{i+1}.png")
        pix.save(image_path)
        image_paths.append(image_path)

    doc.close()
    return image_paths


def save_parsed_output(pages, output_path: str):
    """Save parsed pages to a JSON file."""
    with open(output_path, "w") as f:
        json.dump(pages, f, indent=2)
    return output_path
